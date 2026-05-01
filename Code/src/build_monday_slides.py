"""Append Monday 27th April update slides to the Dr Urman deck.

Casual status update — 5 new slides (11–15). Restored idempotently from
backup.pptx each run. Uses the UZH 'Titel und Inhalt' and 'Nur Titel'
layouts. Image slides and the two-column slide are built with manually
placed shapes so nothing gets cropped by placeholder auto-fit and
columns align properly below the title.

A manual footer strip (Social Computing Group · date · slide number) is
added to every new slide, since python-pptx does not auto-clone the
layout's date/footer/slide-number placeholders onto new slides.
"""

import shutil

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from babelbias.paths import PROJECT_ROOT

DECK_DIR = PROJECT_ROOT / "Presentations" / "27 April"
BACKUP = DECK_DIR / "Monday 27th April.backup.pptx"
DECK = DECK_DIR / "Monday 27th April.pptx"
FIGS = PROJECT_ROOT / "Presentations" / "figures"
MEETING_DATE = "27.04.26"

SLIDE_W = 13.333
SLIDE_H = 7.5
FOOTER_Y = 7.03


# ---------- core helpers ----------

def get_layout(pres, name):
    return next(l for l in pres.slide_layouts if l.name == name)


def ph_by_idx(slide, idx):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    raise KeyError(f"no placeholder idx={idx}")


def _parse_bold(text):
    parts = text.split("**")
    return [(s, i % 2 == 1) for i, s in enumerate(parts)]


def _set_runs(p, segments):
    """Replace p's runs, preserving paragraph-level formatting (level, bullet)."""
    for r in list(p._p.findall(qn("a:r"))):
        p._p.remove(r)
    for text, bold in segments:
        if not text:
            continue
        run = p.add_run()
        run.text = text
        if bold:
            run.font.bold = True


def write_bullets(text_frame, bullets):
    text_frame.word_wrap = True
    for i, b in enumerate(bullets):
        level, text = b if isinstance(b, tuple) else (0, b)
        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        p.level = level
        _set_runs(p, _parse_bold(text))


# ---------- footer strip ----------

def _add_text(slide, left, top, width, height, text, *, size=10,
              bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Inches(0.02)
    tf.margin_top = tf.margin_bottom = Inches(0)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    return tb


def add_footer_strip(slide, slide_num):
    """Match existing deck footer: 'Social Computing Group' after the master's
    'University of Zurich |', and 'DD.MM.YY | N' on the right."""
    _add_text(slide, 1.68, FOOTER_Y, 6.0, 0.22,
              "Social Computing Group", size=9)
    _add_text(slide, 10.80, FOOTER_Y, 2.25, 0.22,
              f"{MEETING_DATE}   |   {slide_num}",
              size=9, align=PP_ALIGN.RIGHT)


# ---------- slide builders ----------

def add_bullet_slide(pres, title, bullets, slide_num):
    slide = pres.slides.add_slide(get_layout(pres, "Titel und Inhalt"))
    ph_by_idx(slide, 0).text = title
    write_bullets(ph_by_idx(slide, 1).text_frame, bullets)
    add_footer_strip(slide, slide_num)
    return slide


def add_title_only(pres, title):
    slide = pres.slides.add_slide(get_layout(pres, "Nur Titel"))
    ph_by_idx(slide, 0).text = title
    return slide


def add_image_and_bullets(pres, title, image_path, bullets, slide_num,
                          img_w=7.5):
    """Image on the left (aspect preserved), bullets on the right."""
    slide = add_title_only(pres, title)

    img_left = 0.35
    img_top = 1.40
    pic = slide.shapes.add_picture(str(image_path),
                                    Inches(img_left), Inches(img_top),
                                    width=Inches(img_w))
    img_h_in = pic.height / 914400
    max_h = SLIDE_H - img_top - 0.70
    if img_h_in > max_h:
        sp = pic._element
        sp.getparent().remove(sp)
        pic = slide.shapes.add_picture(str(image_path),
                                        Inches(img_left), Inches(img_top),
                                        height=Inches(max_h))
        img_w = pic.width / 914400

    # Vertically center the picture within the content band.
    content_top = 1.25
    content_h = SLIDE_H - content_top - 0.70
    pic_h_in = pic.height / 914400
    pic.top = Inches(content_top + (content_h - pic_h_in) / 2)

    text_left = img_left + img_w + 0.40
    text_w = SLIDE_W - text_left - 0.30
    tb = slide.shapes.add_textbox(
        Inches(text_left), Inches(content_top),
        Inches(text_w), Inches(content_h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)

    for i, b in enumerate(bullets):
        level, text = b if isinstance(b, tuple) else (0, b)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level
        _set_runs(p, _parse_bold(text))
        for r in p.runs:
            r.font.size = Pt(20 if level == 0 else 17)
        p.space_after = Pt(8)

    add_footer_strip(slide, slide_num)
    return slide


def add_two_column(pres, title, left_header, left_points,
                    right_header, right_points, slide_num):
    slide = add_title_only(pres, title)

    col_top = 1.40
    col_h = SLIDE_H - col_top - 0.70
    gap = 0.60
    col_w = (SLIDE_W - 0.30 * 2 - gap) / 2
    left_x = 0.30
    right_x = left_x + col_w + gap

    def build_column(x, header, points):
        tb = slide.shapes.add_textbox(
            Inches(x), Inches(col_top), Inches(col_w), Inches(col_h))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.05)

        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = header
        r.font.size = Pt(22)
        r.font.bold = True
        p.space_after = Pt(14)

        for point in points:
            p = tf.add_paragraph()
            _set_runs(p, _parse_bold(point))
            for r in p.runs:
                r.font.size = Pt(18)
            p.space_after = Pt(10)

    build_column(left_x, left_header, left_points)
    build_column(right_x, right_header, right_points)
    add_footer_strip(slide, slide_num)
    return slide


# ---------- deck ----------

def main():
    shutil.copy(BACKUP, DECK)
    pres = Presentation(str(DECK))

    # 11 — what got built this week
    add_bullet_slide(pres, "Prompting pipeline — this week", [
        "9 questions × 3 languages × 10 samples  →  **810 LLM calls**",
        "Three providers: GPT-4o-mini · Claude Haiku 4.5 · Gemini 2.5 Flash",
        "All 810 responses embedded with text-embedding-3-small into the "
        "Wikipedia space",
        "Temperature = 1.0 · resumable · parameterised",
        "Total cost: **~$0.68**",
    ], slide_num=11)

    # 12 — headline result
    add_image_and_bullets(pres,
        "Ingroup bias — confirmed across all three models",
        FIGS / "01_ingroup_bars.png",
        [
            "Every model pulls toward **own-language** Wikipedia",
            (1, "EN → EN wiki:  **+0.18**"),
            (1, "RU → RU wiki:  **+0.05**"),
            (1, "UK → UK wiki:  **+0.04**"),
            "Pattern **identical** across OpenAI, Anthropic, Google",
            "EN pull ~5× stronger than RU/UK — consistent with training-corpus dominance",
        ], slide_num=12, img_w=7.5)

    # 13 — case study
    add_two_column(pres,
        "Case study: 'who are the little green men?'",
        "English prompt", [
            "Defaults to the **sci-fi / extraterrestrial** reading",
            "Mentions Roswell, Area 51, Mars",
            "No mention of Crimea, Russia, or 2014",
        ],
        "Russian / Ukrainian prompts", [
            "Correctly identifies the **2014 Crimea** reference",
            "Unmarked Russian troops, insignia removed",
            "Situates within the annexation timeline",
        ], slide_num=13)

    # 14 — limitations
    add_bullet_slide(pres, "Limitations worth flagging", [
        "'Closer to RU wiki' ≠ pro-Russian content — could be lexical overlap",
        "Wikipedia is in the training data — models may parrot rather than bias",
        "Single event domain so far — need to replicate on a second contested event",
    ], slide_num=14)

    # 15 — next steps
    add_bullet_slide(pres, "Next steps — three parallel tracks", [
        "**Validate** — hand-code stance on a sample of responses to "
        "separate framing bias from lexical similarity",
        "**Generalize** — replicate on a second contested event "
        "(Israel–Palestine · Kashmir · Taiwan)",
        "**Frame** — connect to Oeberst et al. 2020 on Wikipedia ingroup "
        "bias — direct LLM extension",
    ], slide_num=15)

    pres.save(str(DECK))
    print(f"Rebuilt deck: {len(pres.slides)} slides at {DECK}")


if __name__ == "__main__":
    main()
